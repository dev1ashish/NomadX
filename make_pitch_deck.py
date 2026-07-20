#!/usr/bin/env python3
"""Atlas Raman — full pitch deck (clean white background).

Expands the original 6-slide skeleton into a complete narrative, with the
biology and MCR-ALS angles given dedicated, in-depth treatment.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ---------- palette ----------
INK    = RGBColor(0x1A, 0x20, 0x2C)
ACCENT = RGBColor(0x0E, 0x7C, 0x66)   # teal
AMBER  = RGBColor(0xB4, 0x6A, 0x00)   # warm accent for "gotcha"/leakage
MUTE   = RGBColor(0x52, 0x60, 0x6D)
FAINT  = RGBColor(0xEE, 0xF3, 0xF1)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Arial"
MONO   = "Courier New"

HERE = os.path.dirname(os.path.abspath(__file__))
IMG  = os.path.join(HERE, "FINAL", "images")
NB   = os.path.join(HERE, "FINAL", "notebooks")
MX   = Inches(0.7)
BODY_TOP = Inches(1.66)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def new_slide():
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = WHITE
    return s


def no_shadow(sh):
    sh.shadow.inherit = False


def title(slide, text, kicker=None):
    if kicker:
        kb = slide.shapes.add_textbox(MX, Inches(0.34), Inches(11.9), Inches(0.3))
        kp = kb.text_frame.paragraphs[0]
        kr = kp.add_run(); kr.text = kicker.upper()
        kr.font.size = Pt(11.5); kr.font.bold = True; kr.font.name = FONT; kr.font.color.rgb = ACCENT
    tb = slide.shapes.add_textbox(MX, Inches(0.64), Inches(12.0), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(25); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = INK
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MX, Inches(1.34), Inches(2.5), Pt(3))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; rule.line.fill.background(); no_shadow(rule)


def body(slide, specs, left=MX, top=BODY_TOP, width=Inches(11.9), height=Inches(5.3)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for s in specs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(s.get("sb", 0)); p.space_after = Pt(s.get("sa", 6))
        p.line_spacing = s.get("line", 1.03)
        r = p.add_run(); r.text = s.get("prefix", "") + s["text"]
        f = r.font
        f.size = Pt(s.get("size", 15)); f.bold = s.get("bold", False)
        f.italic = s.get("italic", False); f.name = FONT
        f.color.rgb = s.get("color", INK)
    return tb


def b0(t, **k):
    return dict(prefix="•  ", text=t, size=k.get("size", 14.5), bold=k.get("bold", False),
                color=k.get("color", INK), sa=k.get("sa", 7), sb=k.get("sb", 0))


def b1(t, **k):
    return dict(prefix="      –  ", text=t, size=k.get("size", 12), bold=k.get("bold", False),
                color=k.get("color", MUTE), sa=k.get("sa", 3), sb=k.get("sb", 0))


def lead(t, **k):
    return dict(prefix="", text=t, size=k.get("size", 14.5), bold=True,
                color=k.get("color", ACCENT), sa=k.get("sa", 3), sb=k.get("sb", 6))


def note(t, **k):
    return dict(prefix="", text=t, size=k.get("size", 12.5), italic=True,
                color=k.get("color", MUTE), sa=k.get("sa", 0), sb=k.get("sb", 9))


def add_image(slide, path, left, top, max_w, max_h, cap=None):
    if not os.path.exists(path):
        print("  ! missing:", path); return None
    pic = slide.shapes.add_picture(path, left, top, width=max_w)
    if pic.height > max_h:
        pic._element.getparent().remove(pic._element)
        pic = slide.shapes.add_picture(path, left, top, height=max_h)
    if pic.width < max_w:
        pic.left = int(left + (max_w - pic.width) / 2)
    if cap:
        cb = slide.shapes.add_textbox(pic.left, pic.top + pic.height + Pt(3), pic.width, Inches(0.4))
        cp = cb.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cb.text_frame.word_wrap = True
        cr = cp.add_run(); cr.text = cap
        cr.font.size = Pt(9.5); cr.font.italic = True; cr.font.name = FONT; cr.font.color.rgb = MUTE
    return pic


def add_image_center(slide, path, top, max_w, max_h, cap=None):
    return add_image(slide, path, int((prs.slide_width - max_w) / 2), top, max_w, max_h, cap)


def callout(slide, text, left, top, width, height, color=ACCENT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = FAINT
    box.line.color.rgb = color; box.line.width = Pt(1); no_shadow(box)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(12); r.font.italic = True; r.font.name = FONT; r.font.color.rgb = INK
    return box


def footer(slide, n):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MX, Inches(7.0), Inches(11.93), Pt(0.75))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xDD, 0xE3, 0xE1)
    ln.line.fill.background(); no_shadow(ln)
    tb = slide.shapes.add_textbox(Inches(11.7), Inches(7.08), Inches(1.5), Inches(0.34))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(10); r.font.name = FONT; r.font.color.rgb = MUTE
    tb2 = slide.shapes.add_textbox(MX, Inches(7.08), Inches(8.0), Inches(0.34))
    r2 = tb2.text_frame.paragraphs[0].add_run()
    r2.text = "Atlas · Raman bacterial classifier"
    r2.font.size = Pt(10); r2.font.name = FONT; r2.font.color.rgb = MUTE


# =====================================================================
# TITLE
# =====================================================================
s = new_slide()
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MX, Inches(2.5), Inches(2.5), Pt(4))
rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; rule.line.fill.background(); no_shadow(rule)
tb = s.shapes.add_textbox(MX, Inches(2.75), Inches(11.9), Inches(1.7))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = "Classifying Bacteria from Raman Spectra"
r.font.size = Pt(40); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = INK
p2 = tf.add_paragraph(); p2.space_before = Pt(10)
r2 = p2.add_run(); r2.text = "STEC · harmless E. coli · Salmonella · water — a confocal-Raman classifier"
r2.font.size = Pt(19); r2.font.name = FONT; r2.font.color.rgb = MUTE
tb3 = s.shapes.add_textbox(MX, Inches(5.85), Inches(11.9), Inches(0.5))
r3 = tb3.text_frame.paragraphs[0].add_run()
r3.text = "87 samples   •   7,122 spectra   •   9 strains   •   evaluated leave-one-strain-out"
r3.font.size = Pt(14); r3.font.bold = True; r3.font.name = FONT; r3.font.color.rgb = ACCENT
tb4 = s.shapes.add_textbox(MX, Inches(6.3), Inches(11.9), Inches(0.4))
r4 = tb4.text_frame.paragraphs[0].add_run(); r4.text = "NomadX take-home · 2026"
r4.font.size = Pt(12); r4.font.name = FONT; r4.font.color.rgb = MUTE

# =====================================================================
# 2 — THE PROBLEM & THE DATA
# =====================================================================
s = new_slide()
title(s, "The Problem & The Data", "the problem")
body(s, [
    b0("The task: from one bacterium's Raman spectrum, decide which of 4 classes it is — "
       "STEC, harmless E. coli, Salmonella, or water (blank)."),
    b0("Why it matters: STEC (e.g. O157:H7) is the food-poisoning pathogen — telling it from "
       "harmless E. coli fast is real food-safety value."),
    b0("What we got: 87 confocal-Raman maps from one lab (86 .xls + 1 .txt), ~72 pixel-spectra "
       "each, across 9 strains + water."),
    b0("Quality control: 7,999 raw pixels → drop 6 low-signal + 871 background → 7,122 kept (89%). "
       "Each spectrum cropped 2,048 → 987 informative points."),
    b0("Honest headline up front: best classifier ≈ 60% per-class recall on unseen strains "
       "(chance = 25%). Strong where biology allows; near-chance where it doesn't.", bold=True),
], width=Inches(6.1))
add_image(s, os.path.join(IMG, "summary_01_inventory.png"), Inches(6.95), Inches(1.7),
          Inches(6.0), Inches(4.7), cap="87 files / 7,122 spectra · 4 classes · 9 strains + water")
footer(s, 2)

# =====================================================================
# 3 — WHY IT'S HARD: THE ONE-GENE CRUX  (BIOLOGY core)
# =====================================================================
s = new_slide()
title(s, "Why It's Hard — The One-Gene Crux", "why it's hard · biology")
body(s, [
    b0("STEC and harmless E. coli are the same species. What makes one deadly is a single "
       "Shiga-toxin gene, delivered by a bacteriophage (a virus that infects bacteria) — "
       "essentially one functional addition to an ordinary E. coli genome.", bold=True),
    b0("The boundary is virulence-defined, not phylogenetic: the two organisms are nearly "
       "identical; only the acquired toxin gene differs."),
    b0("Raman reads bulk molecular composition — the average protein, lipid, nucleic-acid and "
       "sugar content of the cell. The toxin protein is a tiny fraction of the total."),
    b0("So separating STEC from Non-STEC asks Raman to see the bulk-chemistry shadow of one "
       "gene. Near-chance on this split is the biologically expected result — not a failure of effort.",
       bold=True, color=ACCENT),
    b0("Two hard characters: K-12 — a 100-year-old lab strain with deleted surface genes that no "
       "longer looks like wild E. coli; and Salmonella — a different genus, so genuinely easier."),
], width=Inches(11.9))
callout(s, "Where we DO see between-strain signal, it lives in the cell-wall sugar (LPS) region — "
           "serogroup architecture, not the toxin itself.  →  next slide.",
        MX, Inches(6.05), Inches(11.9), Inches(0.78))
footer(s, 3)

# =====================================================================
# 4 — WHAT THE SPECTRA ARE MADE OF: MACROMOLECULES & LPS  (BIOLOGY detail)
# =====================================================================
s = new_slide()
title(s, "What the Spectra Are Made Of — Macromolecules & LPS", "biology · where the signal lives")
body(s, [
    lead("Every Raman peak belongs to a chemistry family. Five matter here (example peak):"),
    b1("Aromatic amino acids — 1004 cm⁻¹ (phenylalanine): the total-protein anchor"),
    b1("Protein amide / backbone — 1242 (β-sheet), 1658 (α-helix): protein secondary structure"),
    b1("Nucleic acid (DNA/RNA) — 786, 1338"),
    b1("Lipid / carbohydrate — 1080, 1450, 2850–2930: membrane & sugars"),
    b1("LPS (cell-wall sugar) — 1050, 1117, 1194:  the real discriminator", bold=True, color=ACCENT),
    b0("LPS = lipopolysaccharide, the sugar-and-fat coat on the outer membrane. Its O-antigen "
       "sugar chains differ by serogroup (O157 vs O121 vs O103), so the 800–1200 cm⁻¹ LPS region "
       "is where genuine between-strain signal lives.", sb=6),
    b0("Biology-grounded feature confirms it: bio_alpha_helix_score d = −0.986 — Non-STEC carries "
       "more α-helix protein than STEC."),
    b0("The trap: LPS separates serogroups, not virulence → a feature can ace trained strains and "
       "fail on a held-out strain with a different O-antigen.", color=AMBER),
], width=Inches(6.6))
add_image(s, os.path.join(IMG, "05_macromolecule_radar.png"), Inches(7.55), Inches(2.4),
          Inches(5.4), Inches(3.6),
          cap="Macromolecule profile per class (left) and the LPS bands (right): lps_1194 d=+1.03, STEC > Non-STEC.")
footer(s, 4)

# =====================================================================
# 5 — THE PLAN ①: SURVEY + LET THE DATA TALK (EDA)
# =====================================================================
s = new_slide()
title(s, "The Plan ① — Survey the Field, Then Let the Data Talk", "the plan · 1 of 2")
body(s, [
    lead("Step 1 — survey what's known:"),
    b1("Two papers worked this exact problem: Cisek 2013 (STEC vs E. coli) and Yuan 2024 (Salmonella)."),
    b1("We didn't borrow their spectra — Raman is instrument-dependent, so it would teach the machine, not the microbe."),
    lead("Step 2 — let the data talk (exploratory data analysis):", sb=6),
    b1("Class mean spectra are nearly identical → the signal is subtle (shape & ratios, not gross peaks)."),
    b1("Ranking bands by how well they separate classes → the LPS sugar region wins (lps_1194 AUROC 0.775, d = +1.03)."),
    b1("The published Cisek bands came out null / sign-reversed here → don't trust literature bands blindly."),
], width=Inches(11.9), height=Inches(2.6))
add_image_center(s, os.path.join(IMG, "09_best1d_auroc.png"), Inches(4.3),
                 Inches(11.2), Inches(2.25),
                 cap="Best single-band separation per task — LPS bands lead both; everything sits below the 0.75 line.")
footer(s, 5)

# =====================================================================
# 6 — THE PLAN ②: THE HONEST TEST, PRIORITIES & HYPOTHESIS
# =====================================================================
s = new_slide()
title(s, "The Plan ② — The Honest Test, Priorities & Hypothesis", "the plan · 2 of 2")
body(s, [
    b0("The most important decision was not a model — it was the test. Within-strain files share "
       "an acquisition session, so ordinary cross-validation would put near-replicates in train AND "
       "test and leak batch effects, flattering the score.", bold=True),
    b0("So we use Leave-One-Strain-Out (LOSO): train on some strains, test on a strain never "
       "seen. That is the number we trust. (It is a 9-point statistic → wide error bars, stated honestly.)",
       bold=True, color=ACCENT),
    b0("Where to look — priority by macromolecule (from the EDA): carbohydrate / LPS first "
       "(best surface signature) → protein → lipid → nucleic acid last (broadly conserved, so it "
       "looks the same across these bacteria)."),
    b0("Hypothesis: cell-surface sugar (LPS) chemistry carries enough signal to separate "
       "same-species strains — even though the toxin gene itself is invisible to bulk Raman."),
    b0("Success criterion: beat 25% chance under LOSO, and report honestly where it can't."),
], width=Inches(7.1))
add_image(s, os.path.join(NB, "chstretch_pca_scatter.png"), Inches(8.1), Inches(1.9),
          Inches(4.9), Inches(4.6),
          cap="PCA of the C–H region: water (purple) splits off cleanly; bacteria overlap — real but subtle signal.")
footer(s, 6)

# =====================================================================
# 7 — THE PIPELINE
# =====================================================================
s = new_slide()
title(s, "From Messy Files to Trustworthy Data", "methods · data pipeline")
body(s, [
    b0("87 messy instrument files → one clean dataset. The .xls files mislabel their grid size, "
       "mix separators, and include partial scans — all surfaced as metadata, not hidden."),
    b0("Clean every spectrum (5 steps): de-spike cosmic rays → subtract the fluorescence baseline "
       "(arPLS) → smooth (Savitzky–Golay) → crop to informative bands (2,048 → 987) → normalize (SNV)."),
    b0("Predict per pixel, then average a sample's ~72 pixels into one per-file decision — a file "
       "is what you actually want to classify."),
    b0("Built as a contract: blocking quality gates halt the pipeline if the data is bad, so "
       "nothing broken ever reaches the modeling step.", bold=True),
], width=Inches(6.4))
add_image(s, os.path.join(NB, "before_after_spectra.png"), Inches(7.45), Inches(1.9),
          Inches(5.5), Inches(4.6),
          cap="Raw (fluorescent, spiky) → cleaned, normalized 987-point spectrum.")
footer(s, 7)

# =====================================================================
# 8 — MODELING
# =====================================================================
s = new_slide()
title(s, "We Swept the Toolbox — Classical Won", "methods · models")
body(s, [
    b0("Swept the whole toolbox against the honest LOSO bar: PLS-DA, logistic regression, SVM, "
       "random forest, XGBoost (classical); 1D-CNN, 1D-Transformer, domain-adversarial net (deep); "
       "plus 5 ensemble schemes."),
    b0("Winner: PLS-DA on the raw cleaned spectrum — LOSO 0.603. Nothing beat it.", bold=True, color=ACCENT),
    b0("Why classical won: 87 files is tiny. Deep nets overfit the training strains — the "
       "wide-patch Transformer literally blurred the narrow Raman peaks (LOSO 0.19). PLS-DA's "
       "low-rank projection is a built-in regularizer that generalizes to unseen strains."),
    b0("Lesson: for narrow-peak spectroscopy at small N, classical low-rank methods beat flexible "
       "deep models."),
], width=Inches(6.4))
add_image(s, os.path.join(IMG, "fig08_algo_compare.png"), Inches(7.45), Inches(2.0),
          Inches(5.5), Inches(4.2),
          cap="On engineered features LogReg (0.436) beats PLS-DA (0.324) — but raw-spectrum PLS-DA (0.603, dashed) beats them all.")
footer(s, 8)

# =====================================================================
# 9 — FEATURE ENGINEERING
# =====================================================================
s = new_slide()
title(s, "Going Deeper — 259 Engineered Features", "methods · feature engineering")
body(s, [
    lead("To try to beat the raw spectrum, we built 259 chemistry-grounded features in 5 families:"),
    b1("15A · Pseudo-Voigt peak fits (height / width / area per band) — highest-yield; dominates the final model"),
    b1("15B · wavelets, region-PCA, spectral-angle template matching"),
    b1("15C · MCR-ALS unmixing — see next slide"),
    b1("15D · biology ratios — bio_alpha_helix_score d = −0.986 (Non-STEC more α-helix); first K-12-specific axis"),
    b1("15E · spatial moments — null for STEC↔Non-STEC; gives an E. coli↔Salmonella signal (spat_skew_lps_1117 d=+0.725)"),
    b0("Inside each fold, pick the top 35 features by mutual information — with only ~3 features "
       "per file, aggressive selection is mandatory.", sb=6),
    b0("Honest punchline: the best engineered model (LogReg 0.448) still lost to raw-spectrum "
       "PLS-DA (0.603). The information was already in the spectrum shape.", bold=True, color=AMBER),
], width=Inches(6.6))
add_image(s, os.path.join(IMG, "fig02_top_features.png"), Inches(7.55), Inches(2.2),
          Inches(5.4), Inches(4.0),
          cap="Top features by effect size — peak-fits and LPS bands lead.")
footer(s, 9)

# =====================================================================
# 10 — MCR-ALS  (THE CROWN JEWEL + LEAKAGE LESSON)
# =====================================================================
s = new_slide()
title(s, "MCR-ALS — Unmixing, and the Leakage Lesson", "methods · mcr-als")
body(s, [
    b0("MCR-ALS (Multivariate Curve Resolution) un-mixes the whole pile of spectra into a few "
       "“pure-ingredient” spectra plus how much of each sits in every pixel — like reversing a "
       "smoothie back into its ingredients."),
    dict(prefix="     ", text="all spectra  (N pixels × wavenumbers)   ≈   C · Sᵀ      "
         "[ S = K pure-component spectra ;  C = abundance per pixel ]",
         size=12.5, color=INK, sb=2, sa=8),
    b0("We found K = 7 meaningful components — one captured the substrate / fluorescence artifact, "
       "which usefully isolates the instrument signal away from the biology."),
    b0("The crown-jewel lesson: the strongest single feature in the whole project, mcr_C6_mean, "
       "had Cohen's d = −1.23 — but it came from fitting MCR on ALL the data at once (a global fit).",
       bold=True),
    b0("Re-fit MCR INSIDE each fold on training pixels only (the leakage-safe way) and the "
       "components reorder — MCR has rotational / label ambiguity — so “C6” isn't stable and the "
       "feature did NOT survive per-fold selection. The −1.23 was partly a leakage artifact, and "
       "proper cross-validation demoted it.", bold=True, color=AMBER),
], width=Inches(7.6))
add_image(s, os.path.join(IMG, "stage15c_pure_spectra.png"), Inches(8.55), Inches(1.7),
          Inches(4.4), Inches(4.8),
          cap="Global-fit components. Re-fit per fold, these reorder — which is exactly why the d=−1.23 feature isn't reproducible.")
footer(s, 10)

# =====================================================================
# 11 — RESULTS: THE NUMBERS
# =====================================================================
s = new_slide()
title(s, "The Results — The Numbers That Matter", "results")
body(s, [
    b0("The funnel: 87 files → 7,122 spectra → 987 bins → 259 features → 35 selected."),
    dict(prefix="", text="Two recipes — do not mix them up:", size=13.5, bold=True,
         color=INK, sa=2, sb=2),
], width=Inches(6.7), height=Inches(0.95))
# --- recipe comparison card ---
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MX, Inches(2.5), Inches(6.6), Inches(1.42))
panel.fill.solid(); panel.fill.fore_color.rgb = FAINT
panel.line.color.rgb = ACCENT; panel.line.width = Pt(1); no_shadow(panel)
ptf = panel.text_frame; ptf.word_wrap = True
ptf.margin_left = Inches(0.18); ptf.margin_right = Inches(0.12); ptf.margin_top = Inches(0.1)
def _row(p, tag, tagcolor, rest, restbold=True, restcolor=INK):
    r1 = p.add_run(); r1.text = tag
    r1.font.size = Pt(10); r1.font.bold = True; r1.font.name = FONT; r1.font.color.rgb = tagcolor
    r2 = p.add_run(); r2.text = "   " + rest
    r2.font.size = Pt(12.5); r2.font.bold = restbold; r2.font.name = FONT; r2.font.color.rgb = restcolor
p1 = ptf.paragraphs[0]; p1.space_after = Pt(6)
_row(p1, "HEADLINE", ACCENT, "raw 987-bin spectrum → PLS-DA → LOSO 0.603")
p2 = ptf.add_paragraph(); p2.space_after = Pt(6)
_row(p2, "DEPLOYED", MUTE, "35 features → LogReg-L2 → 0.448  (CI 0.345–0.552)", restbold=False)
p3 = ptf.add_paragraph()
r3 = p3.add_run(); r3.text = "PLS-DA on those features collapses to 0.324 — it needs the full spectrum."
r3.font.size = Pt(10.5); r3.font.italic = True; r3.font.name = FONT; r3.font.color.rgb = MUTE
# --- supporting bullets ---
body(s, [
    b0("The engineered track plateaued just below raw — and we report that honestly."),
    b0("Statistics done right: 5,000-sample bootstrap for the error bar; McNemar test "
       "(LogReg beats PLS-DA, p = 0.002 — but only on the same 35 features, NOT a claim against the "
       "0.603 raw headline)."),
    b0("Context: chance is 25%. 0.448 ≈ 1.8× chance, 0.603 ≈ 2.4× chance."),
], width=Inches(6.7), top=Inches(4.15), height=Inches(2.25))
add_image(s, os.path.join(IMG, "fig11_bootstrap_ci.png"), Inches(7.6), Inches(2.1),
          Inches(5.4), Inches(4.2),
          cap="Bootstrap distribution of LOSO accuracy → 95% CI [0.345, 0.552].")
footer(s, 11)

# =====================================================================
# 12 — WHERE IT WORKS & WHERE IT FAILS
# =====================================================================
s = new_slide()
title(s, "Where It Works, and Where It Fails", "results · breakdown")
body(s, [
    b0("The whole-task number hides a split — easy decisions are strong; the biologically-thin one "
       "is near chance."),
    b0("Wins: ATCC25922 0.889, O121:H19 0.889, Typhimurium 0.778. Water-vs-bacteria ≈ AUROC 1.0 in "
       "deployment.", color=ACCENT),
    b0("Failures: K-12 0.00 (atypical lab strain), Dublin 0.111. The held-out water fold scores "
       "0.00 under LOSO — a protocol artifact (no water in training), not a real defect.", color=AMBER),
    b0("The consistent bias: under uncertainty the model defaults to STEC — all 8 held-out water "
       "files were called STEC. Conservative for food safety, but inflates false positives."),
], width=Inches(11.9), height=Inches(2.2))
add_image(s, os.path.join(IMG, "fig06_stage15f_per_strain.png"), Inches(0.8), Inches(4.05),
          Inches(6.0), Inches(2.5), cap="Per-strain LOSO recall — wins and zeros.")
add_image(s, os.path.join(IMG, "fig07_stage15f_confusion.png"), Inches(7.0), Inches(4.05),
          Inches(5.5), Inches(2.5), cap="Confusion matrix — the water row defaults to STEC.")
footer(s, 12)

# =====================================================================
# 13 — THREE SCIENTIFIC FINDINGS
# =====================================================================
s = new_slide()
title(s, "Three Scientific Findings", "what we learned")
body(s, [
    lead("1.  A published discriminator did not replicate."),
    b1("Cisek 2013's three bands (1338 / 1454 / 1658) came out null / sign-reversed on our data "
       "(one band flipped sign). Reported as non-replication under different conditions — not “they're wrong.”"),
    lead("2.  The real signal is cell-wall sugar (LPS).", sb=9),
    b1("Best single band auc_lps_1194, d = +1.03 — surface serogroup architecture in the "
       "800–1200 cm⁻¹ region, exactly where the biology predicted."),
    lead("3.  Our strongest feature was a leakage artifact — and we caught it.", sb=9, color=AMBER),
    b1("mcr_C6_mean (d = −1.23) vanished under per-fold refitting. It is the best evidence that the "
       "whole workflow can be trusted: it demotes its own too-good-to-be-true numbers.", color=AMBER),
], width=Inches(11.9))
footer(s, 13)

# =====================================================================
# 14 — RIGOR
# =====================================================================
s = new_slide()
title(s, "Rigor — How We Avoided Fooling Ourselves", "rigor")
body(s, [
    b0("The real risk in an AI-accelerated workflow isn't bad code — it's confident, plausible, "
       "wrong results. Every step was built to catch exactly that.", bold=True),
    b0("Per-fold refitting of MCR / PCA / SAM — exposes leakage (it killed our own d = −1.23 feature)."),
    b0("Bootstrap confidence intervals, not point estimates. McNemar tests, not eyeballing."),
    b0("Leave-one-strain-out, not flattering k-fold. Pinned-dependency, reproducible notebooks."),
    b0("A literature-falsification track — we tested the Cisek claim instead of assuming it."),
], width=Inches(11.9), height=Inches(3.4))
callout(s, "“Treat every result — AI-generated or not — like a finding from a collaborator you don't "
           "fully trust: verify everything. The leakage I caught in my own strongest feature is the proof.”",
        MX, Inches(5.5), Inches(11.9), Inches(1.0))
footer(s, 14)

# =====================================================================
# 15 — LIMITATIONS & DEPLOYMENT
# =====================================================================
s = new_slide()
title(s, "Honest Limitations & Deployment Behaviour", "limitations")
body(s, [
    b0("STEC ↔ Non-STEC is near chance — biology, not effort (one phage-encoded gene)."),
    b0("LPS separates serogroups, not virulence → a held-out strain with a new O-antigen (e.g. "
       "O157) fails. The “serogroup-specificity trap.”", color=AMBER),
    b0("Tiny N: only 9 strains, one lab, one instrument → wide CI and batch confounding "
       "(a probe still reads file-identity at ~14% vs 1.15% chance)."),
    b0("Mixed-sample stress test: ~10–20% accuracy drop at 10–20% contamination, and the model "
       "leans toward STEC under uncertainty."),
    b0("No cross-instrument or open-set test yet — every number is within our own corpus.", bold=True),
], width=Inches(6.7))
add_image(s, os.path.join(IMG, "stage7_degradation_curves.png"), Inches(7.6), Inches(2.1),
          Inches(5.4), Inches(4.2),
          cap="Synthetic contamination → accuracy degrades; predictions lean STEC.")
footer(s, 15)

# =====================================================================
# 16 — HOW TO IMPROVE
# =====================================================================
s = new_slide()
title(s, "How to Improve It", "next steps")
body(s, [
    b0("Re-measure each specimen multiple times, in different spectral windows, in triplicate — "
       "turns 1 scan into ~12 datasets per sample and controls mechanical / acquisition error.", bold=True),
    b0("Separate lab strains from fresh field isolates so the model learns biology, not lab adaptation."),
    b0("The real bottleneck is strain diversity, not pixel count → add more strains + a second instrument."),
    b0("Free next step: run our model on public ATCC25922 spectra (Zhu 2022, Ho 2019) — the same "
       "strain in another lab. The sharpest cross-lab test, at zero cost.", bold=True, color=ACCENT),
], width=Inches(11.9), top=Inches(1.6), height=Inches(2.5))
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MX, Inches(4.4), Inches(11.9), Inches(2.25))
panel.fill.solid(); panel.fill.fore_color.rgb = FAINT
panel.line.color.rgb = MUTE; panel.line.width = Pt(0.5); no_shadow(panel)
ptf = panel.text_frame; ptf.word_wrap = True
ptf.margin_left = Inches(0.3); ptf.margin_top = Inches(0.16)
diagram = [
    ("Proposed acquisition protocol   (all scans in triplicate → 12 datasets / sample)", True, ACCENT, 13, FONT),
    ("Sample A  ──►  full scan    400–1800 cm⁻¹    ×3", False, INK, 13, MONO),
    ("          ──►  window 1     400–700  cm⁻¹    ×3", False, INK, 13, MONO),
    ("          ──►  window 2     600–1300 cm⁻¹    ×3", False, INK, 13, MONO),
    ("          ──►  window 3    1000–1800 cm⁻¹    ×3", False, INK, 13, MONO),
]
for i, (txt, bold, color, size, fnt) in enumerate(diagram):
    p = ptf.paragraphs[0] if i == 0 else ptf.add_paragraph()
    p.space_after = Pt(3)
    r = p.add_run(); r.text = txt
    r.font.name = fnt; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
footer(s, 16)

# =====================================================================
# 17 — DASHBOARD
# =====================================================================
s = new_slide()
title(s, "An AI-Integrated Dashboard", "product")
body(s, [
    b0("A web dashboard so a non-specialist can use the model: drop in a spectrum → live "
       "identification in seconds."),
    b0("Runs both models side by side — raw-spectrum PLS-DA and engineered-feature LogReg — with "
       "class probabilities for every call."),
    b0("Easy comparison across many parameters in one place:"),
    b1("annotated spectra with the chemistry bands marked  ·  the LPS region"),
    b1("top discriminating features  ·  the 7 MCR components"),
    b1("confusion matrix  ·  an honest uncertainty / limitations panel"),
    b0("Goal: turn the analysis into a tool — identification + interpretation together, with the "
       "model's limits shown openly rather than hidden.", bold=True, sb=4),
    note("Built with a Next.js front-end and serverless (Modal) inference running the real pipeline."),
], width=Inches(11.9))
footer(s, 17)

# =====================================================================
# 18 — APPENDIX A: QUESTIONS — NUMBERS & MODELS
# =====================================================================
def qa_slide(slide_title, kicker, qa, n):
    s = new_slide()
    title(s, slide_title, kicker)
    specs = []
    for q, a in qa:
        specs.append(dict(prefix="Q  ", text=q, size=14, bold=True, color=ACCENT, sa=1, sb=8))
        specs.append(dict(prefix="      ", text=a, size=12.5, color=INK, sa=1, line=1.04))
    body(s, specs, width=Inches(11.9))
    footer(s, n)


qa_slide("Questions to Expect — Numbers & Models", "appendix · 1 of 2", [
    ("“You have two different best numbers?”",
     "0.603 = PLS-DA on the raw spectrum (headline). 0.448 = LogReg on engineered features (deployed). "
     "The engineered features plateaued below raw — and we say so."),
    ("“LogReg beats PLS-DA, p = 0.002 — so you beat your baseline?”",
     "No. That test is engineered-vs-engineered features only. Raw-spectrum PLS-DA (0.603) was never beaten."),
    ("“Isn't the LogReg 0.488?”",
     "≈ 0.448, not 0.488 — the best model of the feature-engineering track. It still lost to raw-spectrum PLS-DA (0.603)."),
    ("“Same model for water and Salmonella?”",
     "One multiclass model with four outputs — a spectrum scores all four classes and we take the highest. Not one model per class."),
    ("“Why not an ensemble?”",
     "We tried five ensemble schemes plus the deep nets. None beat the single PLS-DA — the added complexity didn't improve generalization."),
    ("“What is PLS-DA, in one breath?”",
     "It squashes the 987-point spectrum into ~30 label-aware directions that best separate the four classes, then classifies there — like PCA, but supervised."),
], 18)

# =====================================================================
# 19 — APPENDIX B: QUESTIONS — METHODS & RIGOR
# =====================================================================
qa_slide("Questions to Expect — Methods & Rigor", "appendix · 2 of 2", [
    ("“Why did deep learning lose to a classical method?”",
     "87 files is tiny — CNNs / Transformers overfit, and wide patches blur the narrow Raman peaks. At small N, classical low-rank wins."),
    ("“Did you really falsify Cisek 2013?”",
     "Careful wording: it didn't transfer to file level under our replicate structure — not “they are wrong.”"),
    ("“Is your strongest feature real?”",
     "No — mcr_C6_mean d = −1.23 came from a global MCR fit (leakage). It was killed by re-fitting MCR inside each fold."),
    ("“Why not generate synthetic data?”",
     "The bottleneck is real strains, not spectra — synthetic data can't invent an unseen strain, would inflate LOSO, and risks teaching artifacts. Even a WGAN tops out ~94%."),
    ("“Why not use the reference papers' data?”",
     "Raman is instrument-dependent — another lab's spectra teach the instrument, not the microbe (a batch confound). We used their methods; running our model on public ATCC25922 is the free #1 next step."),
    ("“Did you actually have cosmic rays?”",
     "Yes — 99.7% of spectra had ≥1 flagged point, the tallest ~6,000σ. But the z>5 threshold also trims ordinary noise; the genuine hits are the extreme tail above z>50."),
], 19)

out = os.path.join(HERE, "Atlas_Raman_Pitch_Deck.pptx")
prs.save(out)
print("Saved:", out, "(%d slides)" % len(prs.slides._sldIdLst))
